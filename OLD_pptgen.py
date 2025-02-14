import streamlit as st
import openai
import json
import os
import torch
from transformers import CLIPProcessor, CLIPModel
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from PIL import Image

# Set OpenAI API key from Streamlit secrets
openai.api_key = st.secrets["openai_api_key"]

# Load CLIP model and processor
device = "cuda" if torch.cuda.is_available() else "cpu"
@st.cache_resource
def load_clip_model():
    model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
    processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
    model.to(device)
    return model, processor

clip_model, clip_processor = load_clip_model()

def generate_slides(user_instructions):
    """
    Uses GPT-4o to generate slide content in JSON format.
    Expected JSON format:
    {
      "slides": [
         {
           "title": "Slide Title",
           "content": "Full text content rewritten based on user instructions.",
           "keywords": ["keyword1", "keyword2"]
         },
         ...
      ]
    }
    """
    prompt = f"""
You are an assistant that completely rewrites presentation content based solely on the user instructions provided.
Do not include any default template text.
Your output should be a JSON object with a single key "slides" containing an array of slide objects.
Each slide object must have:
 - "title": a concise slide title.
 - "content": full text content that completely reflects the user instructions.
 - "keywords": an array of keywords for image selection.

User instructions: {user_instructions}

Please output only the JSON with no additional commentary or markdown formatting.
Example output:
{{
  "slides": [
    {{
      "title": "Introduction",
      "content": "This slide introduces the topic with rewritten content based on the user instructions.",
      "keywords": ["introduction", "overview"]
    }},
    {{
      "title": "Details",
      "content": "This slide provides detailed, rewritten content reflecting the user input.",
      "keywords": ["details", "information"]
    }}
  ]
}}
"""
    try:
        response = openai.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You generate presentation slide content in JSON based solely on user instructions."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
        )
        response_content = response.choices[0].message.content.strip()
        st.write("Raw GPT Output:", response_content)  # Debug output

        # Remove markdown fences if present.
        if response_content.startswith("```"):
            lines = response_content.splitlines()
            if len(lines) > 2:
                response_content = "\n".join(lines[1:-1]).strip()
            else:
                response_content = response_content.strip("```").strip()

        slides_data = json.loads(response_content)
        return slides_data.get("slides", [])
    except Exception as e:
        st.error(f"Error generating slides: {e}")
        return None

def select_best_image_for_slide(slide_data, images_folder="images"):
    """
    Uses CLIP to compute the similarity between the slide's combined text (title, content, keywords)
    and each image in the images folder. Returns the image path with the highest cosine similarity.
    """
    if not os.path.isdir(images_folder):
        return None

    slide_text = " ".join([
        slide_data.get("title", ""),
        slide_data.get("content", ""),
        " ".join(slide_data.get("keywords", []))
    ]).strip()
    if not slide_text:
        return None

    text_inputs = clip_processor(text=[slide_text], return_tensors="pt", padding=True)
    with torch.no_grad():
        text_features = clip_model.get_text_features(**text_inputs).to(device)
    text_features = text_features / text_features.norm(dim=-1, keepdim=True)

    best_score = -1.0
    best_image = None

    for filename in os.listdir(images_folder):
        if filename.lower().endswith((".jpg", ".jpeg", ".png")):
            image_path = os.path.join(images_folder, filename)
            try:
                image = Image.open(image_path).convert("RGB")
            except Exception as e:
                st.write(f"Error opening image {filename}: {e}")
                continue
            image_inputs = clip_processor(images=image, return_tensors="pt")
            with torch.no_grad():
                image_features = clip_model.get_image_features(**image_inputs).to(device)
            image_features = image_features / image_features.norm(dim=-1, keepdim=True)
            score = torch.cosine_similarity(text_features, image_features).item()
            st.write(f"Similarity score for '{filename}': {score:.2f}")
            if score > best_score:
                best_score = score
                best_image = image_path

    st.write("Selected best image:", best_image, "with score:", best_score)
    return best_image

def create_final_presentation(slides, template_path, presentation_title, institution):
    """
    Builds the final presentation using your three-slide master template (template.pptx).
    The template's layouts are assumed to be:
      - Layout 0: Title slide.
      - Layout 1: Content slide.
      - Layout 2: Thank You slide.
    The final presentation will include:
      1. A Title slide (layout 0) updated with presentation_title and institution.
      2. One content slide per GPT-generated slide (using layout 1) with text on the left and an image on the right.
      3. A Thank You slide (layout 2) added as is.
    """
    # Load the master template.
    master_pres = Presentation(template_path)
    # Retrieve the custom layouts.
    title_layout = master_pres.slide_layouts[0]
    content_layout = master_pres.slide_layouts[1]
    thankyou_layout = master_pres.slide_layouts[2]

    # Create a new presentation from the master template.
    final_pres = Presentation(template_path)
    # Remove any default slides using the XML method.
    xml_slides = final_pres.slides._sldIdLst
    for sld in list(xml_slides):
        xml_slides.remove(sld)

    # Determine slide width in inches.
    slide_width_inches = final_pres.slide_width / 914400

    # --- Title Slide ---
    title_slide = final_pres.slides.add_slide(title_layout)
    if title_slide.shapes.title:
        title_slide.shapes.title.text = presentation_title
        title_slide.shapes.title.text_frame.word_wrap = True
    try:
        subtitle = title_slide.placeholders[1]
        subtitle.text = institution
        subtitle.text_frame.word_wrap = True
    except Exception:
        tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(slide_width_inches - 2), Inches(1))
        tb.text_frame.text = institution
        tb.text_frame.word_wrap = True

    # --- Content Slides ---
    # Define layout: reserve 60% for text (left) and 40% for image (right).
    left_margin = 0.5
    right_margin = 0.5
    gap = 0.5
    available_width = slide_width_inches - left_margin - right_margin
    text_width = (available_width * 0.6) - (gap / 2)
    image_width = (available_width * 0.4) - (gap / 2)
    text_left = left_margin
    image_left = left_margin + text_width + gap

    for slide_data in slides:
        slide = final_pres.slides.add_slide(content_layout)
        try:
            slide.shapes.title.text = slide_data.get("title", "")
            slide.shapes.title.text_frame.word_wrap = True
        except Exception:
            # Fallback: add title textbox.
            tbox = slide.shapes.add_textbox(Inches(text_left), Inches(0.5), Inches(text_width), Inches(1))
            tbox.text_frame.text = slide_data.get("title", "")
            tbox.text_frame.word_wrap = True

        try:
            content_ph = slide.placeholders[1]
            content_ph.text = slide_data.get("content", "")
            content_ph.text_frame.word_wrap = True
            content_ph.left = Inches(text_left)
            content_ph.width = Inches(text_width)
        except Exception:
            cbox = slide.shapes.add_textbox(Inches(text_left), Inches(1.8), Inches(text_width), Inches(3))
            cbox.text_frame.text = slide_data.get("content", "")
            cbox.text_frame.word_wrap = True

        best_image = select_best_image_for_slide(slide_data)
        if best_image and os.path.exists(best_image):
            slide.shapes.add_picture(best_image, Inches(image_left), Inches(1), width=Inches(image_width))

    # --- Thank You Slide ---
    thankyou_slide = final_pres.slides.add_slide(thankyou_layout)
    # (Assumes the Thank You slide is already designed as desired in the template.)

    return final_pres

def main():
    st.title("AI-Powered PowerPoint Presentation Generator")
    st.write(
        "Enter your presentation instructions below. The final presentation will include a Title slide (with your "
        "presentation title and university/college), content slides generated from your input (with text on the left and a "
        "best-fit image on the right), and a Thank You slide at the end. All slides will use the theme defined in your master template."
    )

    presentation_title = st.text_input("Presentation Title", "My Presentation")
    institution = st.text_input("University/College", "ABC College")

    # Use the master template file (which now contains 3 slides: Title, Content, Thank You).
    template_path = os.path.join("templates", "template.pptx")
    if not os.path.exists(template_path):
        st.error("Master template file not found in the templates folder!")
        st.stop()

    user_instructions = st.text_area("Presentation Instructions", height=150)

    if st.button("Generate Presentation"):
        if not user_instructions.strip():
            st.warning("Please enter your presentation instructions.")
            return

        with st.spinner("Generating slide content..."):
            slides = generate_slides(user_instructions)

        if slides is None or len(slides) == 0:
            st.error("Failed to generate slide content. Check your instructions and try again.")
            return

        st.subheader("Generated Slide Data")
        st.json(slides)

        with st.spinner("Building final presentation..."):
            final_pres = create_final_presentation(slides, template_path, presentation_title, institution)
            if final_pres is None:
                st.error("Error creating final presentation.")
                return

            pptx_io = BytesIO()
            final_pres.save(pptx_io)
            pptx_io.seek(0)

        st.success("Presentation created successfully!")
        st.download_button(
            label="Download Presentation",
            data=pptx_io,
            file_name="generated_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

if __name__ == "__main__":
    main()
