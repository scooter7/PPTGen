import streamlit as st
import openai
import json
import os
import torch
from transformers import CLIPProcessor, CLIPModel
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from PIL import Image
from copy import deepcopy

# Set OpenAI API key from Streamlit secrets
openai.api_key = st.secrets["openai_api_key"]

# Load CLIP model and processor (this may take a few seconds on first load)
device = "cuda" if torch.cuda.is_available() else "cpu"
@st.cache_resource
def load_clip_model():
    model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
    processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
    model.to(device)
    return model, processor

clip_model, clip_processor = load_clip_model()

def generate_slides(user_instructions):
    """
    Use GPT-4o to generate slide content in JSON format.
    Expected output:
    {
      "slides": [
         {
           "title": "Slide Title",
           "content": "Full text content rewritten to reflect the user instructions.",
           "keywords": ["keyword1", "keyword2"]
         },
         ...
      ]
    }
    """
    prompt = f"""
You are an assistant that completely rewrites presentation content based solely on the user instructions provided.
Do not include any default template text.
Your output should be a JSON object with a single key "slides" containing an array of slide objects.
Each slide object must have:
 - "title": a concise slide title.
 - "content": full text content that completely reflects the user instructions.
 - "keywords": an array of keywords for image selection.

User instructions: {user_instructions}

Please output only the JSON with no additional commentary or markdown formatting.
Example output:
{{
  "slides": [
    {{
      "title": "Introduction",
      "content": "This slide introduces the topic with rewritten content based on the user instructions.",
      "keywords": ["introduction", "overview"]
    }},
    {{
      "title": "Details",
      "content": "This slide provides detailed, rewritten content reflecting the user input.",
      "keywords": ["details", "information"]
    }}
  ]
}}
"""
    try:
        # Use the correct syntax for the chat completions call.
        response = openai.chat.completions.create(
            model="gpt-4o",  # Using GPT-4o as requested
            messages=[
                {"role": "system", "content": "You generate presentation slide content in JSON based solely on user instructions."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
        )
        response_content = response.choices[0].message.content.strip()
        st.write("Raw GPT Output:", response_content)  # Debug output

        # Remove markdown code fences if present.
        if response_content.startswith("```"):
            lines = response_content.splitlines()
            if len(lines) > 2:
                response_content = "\n".join(lines[1:-1]).strip()
            else:
                response_content = response_content.strip("```").strip()

        slides_data = json.loads(response_content)
        return slides_data.get("slides", [])
    except Exception as e:
        st.error(f"Error generating slides: {e}")
        return None

def select_best_image_for_slide(slide_data, images_folder="images"):
    """
    Use CLIP to compute the similarity between the slide's text and each image.
    The slide's text is computed from its title, content, and keywords.
    Returns the image path with the highest cosine similarity.
    """
    if not os.path.isdir(images_folder):
        return None

    # Combine slide text.
    slide_text = " ".join([
        slide_data.get("title", ""),
        slide_data.get("content", ""),
        " ".join(slide_data.get("keywords", []))
    ]).strip()
    if not slide_text:
        return None

    # Compute text embedding.
    text_inputs = clip_processor(text=[slide_text], return_tensors="pt", padding=True)
    with torch.no_grad():
        text_features = clip_model.get_text_features(**text_inputs).to(device)
    text_features = text_features / text_features.norm(dim=-1, keepdim=True)

    best_score = -1.0
    best_image = None

    # Iterate over images.
    for filename in os.listdir(images_folder):
        if filename.lower().endswith((".jpg", ".jpeg", ".png")):
            image_path = os.path.join(images_folder, filename)
            try:
                image = Image.open(image_path).convert("RGB")
            except Exception as e:
                st.write(f"Error opening image {filename}: {e}")
                continue
            image_inputs = clip_processor(images=image, return_tensors="pt")
            with torch.no_grad():
                image_features = clip_model.get_image_features(**image_inputs).to(device)
            image_features = image_features / image_features.norm(dim=-1, keepdim=True)
            score = torch.cosine_similarity(text_features, image_features).item()
            st.write(f"Similarity score for '{filename}': {score:.2f}")
            if score > best_score:
                best_score = score
                best_image = image_path

    st.write("Selected best image:", best_image, "with score:", best_score)
    return best_image

def create_content_presentation(slides, content_template_path):
    """
    Create a presentation containing content slides using the provided content template.
    This function creates slides based on the base slide design in the content template,
    laying out a text area (left) and an image area (right) so that they do not overlap.
    """
    prs = Presentation(content_template_path)
    if len(prs.slides) == 0:
        st.error("Content template does not contain any slides.")
        return None

    # Compute layout parameters.
    slide_width_inches = prs.slide_width / 914400  # EMU to inches
    left_margin = 0.5
    right_margin = 0.5
    gap = 0.5  # gap between text and image areas
    available_width = slide_width_inches - left_margin - right_margin
    text_width = (available_width * 0.6) - (gap / 2)
    image_width = (available_width * 0.4) - (gap / 2)
    text_left = left_margin
    image_left = left_margin + text_width + gap

    st.write(f"Slide width: {slide_width_inches:.2f} inches, Text area: {text_width:.2f} inches, Image area: {image_width:.2f} inches")

    # Use the base slide from the content template.
    base_slide = prs.slides[0]
    first_slide_data = slides[0]

    # Update the base slide's title.
    if base_slide.shapes.title:
        base_slide.shapes.title.text = first_slide_data.get("title", "")
        base_slide.shapes.title.text_frame.word_wrap = True

    # Update (or add) the content textbox.
    try:
        content_placeholder = base_slide.placeholders[1]
        content_placeholder.text = first_slide_data.get("content", "")
        content_placeholder.text_frame.word_wrap = True
        content_placeholder.left = Inches(text_left)
        content_placeholder.width = Inches(text_width)
    except (IndexError, KeyError):
        left = Inches(text_left)
        top = Inches(2)
        width = Inches(text_width)
        height = Inches(3)
        textbox = base_slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.text = first_slide_data.get("content", "")
        textbox.text_frame.word_wrap = True

    # Insert best-fit image.
    best_image = select_best_image_for_slide(first_slide_data)
    if best_image and os.path.exists(best_image):
        base_slide.shapes.add_picture(best_image, Inches(image_left), Inches(1), width=Inches(image_width))

    # Retrieve the base slide's layout.
    base_layout = base_slide.slide_layout

    # Create additional content slides.
    for slide_data in slides[1:]:
        new_slide = prs.slides.add_slide(base_layout)
        if new_slide.shapes.title:
            new_slide.shapes.title.text = slide_data.get("title", "")
            new_slide.shapes.title.text_frame.word_wrap = True
        try:
            content_placeholder = new_slide.placeholders[1]
            content_placeholder.text = slide_data.get("content", "")
            content_placeholder.text_frame.word_wrap = True
            content_placeholder.left = Inches(text_left)
            content_placeholder.width = Inches(text_width)
        except (IndexError, KeyError):
            left = Inches(text_left)
            top = Inches(2)
            width = Inches(text_width)
            height = Inches(3)
            textbox = new_slide.shapes.add_textbox(left, top, width, height)
            textbox.text_frame.text = slide_data.get("content", "")
            textbox.text_frame.word_wrap = True

        best_image = select_best_image_for_slide(slide_data)
        if best_image and os.path.exists(best_image):
            new_slide.shapes.add_picture(best_image, Inches(image_left), Inches(1), width=Inches(image_width))
    return prs

def clone_slide(target_pres, source_slide):
    """
    Clone a slide from a source presentation into the target presentation.
    This hack uses deepcopy on the slide's XML.
    """
    blank_layout = target_pres.slide_layouts[6]  # Use a blank layout
    new_slide = target_pres.slides.add_slide(blank_layout)
    for shape in source_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(deepcopy(shape.element), 'p:extLst')
    return new_slide

def build_final_presentation(content_pres, title_path, thankyou_path, presentation_title, institution):
    """
    Build the final presentation by:
      1. Loading the Title slide from Title.pptx and updating it with the presentation title and institution.
      2. Cloning all content slides from the generated content presentation.
      3. Loading the Thank You slide from ThankYou.pptx and appending it.
    All slides are cloned into a new blank presentation.
    """
    final_pres = Presentation()  # new blank presentation

    # Load and update the Title slide.
    title_pres = Presentation(title_path)
    title_slide = title_pres.slides[0]
    # Update title text (assume the main title placeholder is present).
    if title_slide.shapes.title:
        title_slide.shapes.title.text = presentation_title
        title_slide.shapes.title.text_frame.word_wrap = True
    # Try to update a subtitle placeholder with the institution.
    try:
        subtitle = title_slide.placeholders[1]
        subtitle.text = institution
        subtitle.text_frame.word_wrap = True
    except (IndexError, KeyError):
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(1)
        textbox = title_slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.text = institution
        textbox.text_frame.word_wrap = True

    # Clone the Title slide into the final presentation.
    clone_slide(final_pres, title_slide)

    # Clone each content slide from the generated content presentation.
    for slide in content_pres.slides:
        clone_slide(final_pres, slide)

    # Load the Thank You slide.
    thankyou_pres = Presentation(thankyou_path)
    thankyou_slide = thankyou_pres.slides[0]
    # Clone the Thank You slide into the final presentation.
    clone_slide(final_pres, thankyou_slide)

    return final_pres

def main():
    st.title("AI-Powered PowerPoint Presentation Generator")
    st.write(
        "Enter your presentation instructions below. The output presentation will include a custom Title slide, "
        "content slides generated from your input (with images selected via computer vision), and a Thank You slide at the end."
    )

    # Additional fields for title slide information.
    presentation_title = st.text_input("Presentation Title", "My Presentation")
    institution = st.text_input("University/College", "ABC College")

    # Use fixed paths for title and Thank You slides.
    title_path = os.path.join("templates", "Title.pptx")
    thankyou_path = os.path.join("templates", "ThankYou.pptx")
    # For content slides, use the content template.
    content_template_path = os.path.join("templates", "powerpointtemplate.pptx")

    user_instructions = st.text_area("Presentation Instructions", height=150)

    if st.button("Generate Presentation"):
        if not user_instructions.strip():
            st.warning("Please enter your presentation instructions.")
            return

        with st.spinner("Generating slide content..."):
            slides = generate_slides(user_instructions)

        if slides is None or len(slides) == 0:
            st.error("Failed to generate slide content. Check your instructions and try again.")
            return

        st.subheader("Generated Slide Data")
        st.json(slides)

        with st.spinner("Creating content slides..."):
            content_pres = create_content_presentation(slides, content_template_path)
            if content_pres is None:
                st.error("Error creating content slides.")
                return

        with st.spinner("Building final presentation..."):
            final_pres = build_final_presentation(content_pres, title_path, thankyou_path, presentation_title, institution)
            if final_pres is None:
                st.error("Error building final presentation.")
                return

            pptx_io = BytesIO()
            final_pres.save(pptx_io)
            pptx_io.seek(0)

        st.success("Presentation created successfully!")
        st.download_button(
            label="Download Presentation",
            data=pptx_io,
            file_name="generated_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

if __name__ == "__main__":
    main()
